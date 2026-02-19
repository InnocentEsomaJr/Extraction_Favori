import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import requests
from requests.adapters import HTTPAdapter
from urllib3.util.retry import Retry
import io
import os
from datetime import datetime

# --- 1. CONFIGURATION ---
st.set_page_config(page_title="SNIS RDC - Dashboard Performance", layout="wide")

# Fonctions de coloration (Scripts originaux)
def style_taux(val):
    try: val = float(val)
    except: return ''
    if val < 50: return 'color: white; background-color: #FF0000'
    elif 50 <= val <= 69: return 'color: white; background-color: #800000'
    elif 70 <= val <= 79: return 'color: black; background-color: #FFFF00'
    elif 80 <= val <= 95: return 'color: black; background-color: #32CD32'
    elif val > 95: return 'color: white; background-color: #008000'
    return ''

def style_score(val):
    try: val = int(val)
    except: return ''
    if val < 5: return 'color: white; background-color: #800000'
    elif val == 5: return 'color: black; background-color: #FFC0CB'
    elif 6 <= val <= 9: return 'color: black; background-color: #32CD32'
    elif val == 10: return 'color: white; background-color: #008000'
    elif val > 10: return 'color: white; background-color: #004d00'
    return ''

def normalize_org_name(name):
    """Normalise les noms d'unités pour les comparaisons robustes."""
    if pd.isna(name):
        return ""
    return " ".join(str(name).strip().lower().split())

def _find_best_column(columns, exact_names=None, include_terms=None, exclude_terms=None):
    """
    Trouve la meilleure colonne selon noms exacts ou termes inclus/exclus.
    Retourne None si non trouvée.
    """
    cols = list(columns)
    exact_names = exact_names or []
    include_terms = include_terms or []
    exclude_terms = exclude_terms or []

    for name in exact_names:
        if name in cols:
            return name

    for col in cols:
        col_l = str(col).strip().lower()
        if include_terms and not all(term in col_l for term in include_terms):
            continue
        if exclude_terms and any(term in col_l for term in exclude_terms):
            continue
        return col
    return None

def normalize_orgunit_columns(df):
    """
    Normalise les colonnes organisationnelles vers:
    - Organisation unit
    - Organisation unit ID (si disponible)
    - Organisation unit is parent (si disponible)
    """
    if df is None or not isinstance(df, pd.DataFrame):
        return df

    out = df.copy()
    cols = out.columns.tolist()

    # Colonne nom d'unité
    org_col = _find_best_column(
        cols,
        exact_names=["Organisation unit", "Organisation Unit", "Org unit", "Org Unit"],
        include_terms=["organisation", "unit"],
        exclude_terms=["id", "uid", "code", "description", "parameter", "is parent"]
    )
    if org_col is None:
        org_col = _find_best_column(
            cols,
            include_terms=["org", "unit"],
            exclude_terms=["id", "uid", "code", "description", "parameter", "is parent"]
        )
    if org_col is None and len(cols) > 0:
        # Fallback: première colonne non métrique
        metric_terms = ["reporting rate", "actual reports", "expected reports", "ratio", "score"]
        fallback_cols = [c for c in cols if not any(t in str(c).lower() for t in metric_terms)]
        org_col = fallback_cols[0] if fallback_cols else cols[0]

    # Colonne ID d'unité
    id_col = _find_best_column(
        cols,
        exact_names=["Organisation unit ID", "Organisation Unit ID", "Org unit ID", "Organisation unit UID", "ou"],
        include_terms=["organisation", "unit", "id"]
    )
    if id_col is None:
        id_col = _find_best_column(
            cols,
            include_terms=["org", "unit", "id"]
        )
    if id_col is None:
        id_col = _find_best_column(
            cols,
            include_terms=["organisation", "unit", "uid"]
        )
    if id_col is None:
        id_col = _find_best_column(
            cols,
            include_terms=["org", "unit", "uid"]
        )
    if id_col is None and "ou" in cols:
        id_col = "ou"

    # Colonne parent (optionnelle)
    parent_col = _find_best_column(
        cols,
        exact_names=["Organisation unit is parent"],
        include_terms=["is parent"]
    )
    if parent_col is None:
        parent_col = _find_best_column(
            cols,
            include_terms=["organisation", "unit", "parent"]
        )

    rename_map = {}
    if org_col and org_col != "Organisation unit":
        rename_map[org_col] = "Organisation unit"
    if id_col and id_col != "Organisation unit ID":
        rename_map[id_col] = "Organisation unit ID"
    if parent_col and parent_col != "Organisation unit is parent":
        rename_map[parent_col] = "Organisation unit is parent"
    if rename_map:
        out = out.rename(columns=rename_map)

    # Garantir l'existence de la colonne organisation unit pour éviter les KeyError
    if "Organisation unit" not in out.columns:
        out["Organisation unit"] = "Unité non renseignée"

    return out

def _extract_identifier(value):
    """Extrait un identifiant stable depuis un champ DHIS2 (dict ou string)."""
    if isinstance(value, dict):
        return str(
            value.get("id")
            or value.get("uid")
            or value.get("displayName")
            or value.get("name")
            or ""
        ).strip()
    if value is None:
        return ""
    return str(value).strip()

def build_violation_signature_set(df_val):
    """
    Construit un ensemble de signatures de violations pour comparer T-1 vs T.
    On exclut volontairement la période pour détecter les erreurs réellement corrigées.
    """
    if df_val is None or df_val.empty:
        return set()

    signatures = set()
    for row in df_val.to_dict("records"):
        ou = _extract_identifier(row.get("organisationUnit"))
        rule = _extract_identifier(row.get("validationRule"))
        aoc = _extract_identifier(row.get("attributeOptionCombo"))
        coc = _extract_identifier(row.get("categoryOptionCombo"))
        ds = _extract_identifier(row.get("dataSet"))

        # Une violation doit au minimum avoir une règle ou une OU.
        if not ou and not rule:
            continue

        signatures.add((ou, rule, aoc, coc, ds))
    return signatures

# Fonction pour l'export Excel
def to_excel(df):
    output = io.BytesIO()
    # Utilisation de xlsxwriter (assurez-vous qu'il est installé via pip install xlsxwriter)
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        df.to_excel(writer, index=False, sheet_name='Performance')
    return output.getvalue()

def _build_styler(
    df,
    taux_cols=None,
    score_cols=None,
    int_cols=None,
    decimals=2,
    highlight_max_cols=None,
    ratio_alert_col=None,
    ratio_alert_threshold=10.0
):
    """Construit un Styler pandas avec les mêmes règles de couleurs que le dashboard."""
    if df is None or df.empty:
        return None

    taux_cols = [c for c in (taux_cols or []) if c in df.columns]
    score_cols = [c for c in (score_cols or []) if c in df.columns]
    int_cols = [c for c in (int_cols or []) if c in df.columns]
    highlight_max_cols = [c for c in (highlight_max_cols or []) if c in df.columns]

    styler = df.style

    numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
    format_map = {c: f"{{:.{int(decimals)}f}}" for c in numeric_cols}
    for c in int_cols:
        format_map[c] = "{:.0f}"
    if format_map:
        styler = styler.format(format_map, na_rep="")

    if taux_cols:
        styler = styler.map(style_taux, subset=taux_cols)
    if score_cols:
        styler = styler.map(style_score, subset=score_cols)
    if highlight_max_cols:
        styler = styler.highlight_max(subset=highlight_max_cols, color="#ffcccc")
    if ratio_alert_col and ratio_alert_col in df.columns:
        styler = styler.map(
            lambda x: (
                "color: red; font-weight: bold"
                if pd.notna(x) and pd.to_numeric(x, errors="coerce") > ratio_alert_threshold
                else ""
            ),
            subset=[ratio_alert_col]
        )

    return styler

def _extract_css_values(css_text):
    """Extrait couleur texte/fond/gras depuis une chaîne CSS."""
    default_text = "#000000"
    default_bg = "#FFFFFF"
    bold = False
    if not css_text:
        return default_text, default_bg, bold

    text_color = default_text
    bg_color = default_bg
    for item in str(css_text).split(";"):
        if ":" not in item:
            continue
        key, value = item.split(":", 1)
        key = key.strip().lower()
        value = value.strip()
        if key == "color" and value:
            text_color = value
        elif key == "background-color" and value:
            bg_color = value
        elif key == "font-weight" and "bold" in value.lower():
            bold = True
    return text_color, bg_color, bold

def build_dashboard_comments_df(
    df_final_c,
    df_final_p,
    df_synth,
    selected_zone,
    df_tab4_fusion=None,
    df_tab5=None,
    df_base=None,
    df_actual_detail=None,
    df_expected_detail=None,
    df_comparatif=None,
    df_top5=None,
    df_flop5=None,
):
    """Prépare des commentaires clairs pour graphiques et tableaux exportés."""
    zone_label = selected_zone if selected_zone != "Toutes les zones" else "Toutes les zones"
    comp_avg = float(df_synth['Complétude_Globale (%)'].mean()) if not df_synth.empty else 0.0
    prom_avg = float(df_synth['Promptitude_Globale (%)'].mean()) if not df_synth.empty else 0.0
    top_comp = df_synth.nlargest(1, 'Complétude_Globale (%)')['Organisation unit'].iloc[0] if not df_synth.empty else "N/A"
    flop_prom = df_synth.nsmallest(1, 'Promptitude_Globale (%)')['Organisation unit'].iloc[0] if not df_synth.empty else "N/A"

    rows = [
        {
            "Element": "Légende des couleurs",
            "Commentaire": (
                "Taux: <50 rouge vif, 50-69 bordeaux, 70-79 jaune, 80-95 vert clair, >95 vert foncé. "
                "Score datasets: <5 bordeaux, 5 rose, 6-9 vert clair, 10 vert, >10 vert foncé."
            )
        },
        {
            "Element": "Onglet 2 - Tableau Complétude",
            "Commentaire": (
                f"Complétude moyenne {comp_avg:.2f}%. Les cellules colorées reprennent strictement la logique du dashboard. "
                "Priorité: unités en rouge/bordeaux."
            )
        },
        {
            "Element": "Base de données",
            "Commentaire": (
                "Cette feuille reprend les données sources filtrées (niveau, période, zone). "
                "Elle sert de référence brute pour toutes les analyses."
            )
        },
        {
            "Element": "Rapports détaillés - Données réelles",
            "Commentaire": (
                "Ce tableau détaille les 'Actual reports' par indicateur et la somme Reports_Actual. "
                "Il reflète le volume réel des notifications."
            )
        },
        {
            "Element": "Rapports détaillés - Participants/Attendus",
            "Commentaire": (
                "Ce tableau présente les 'Expected reports' et la somme Reports_Attendu. "
                "Il représente le dénominateur attendu pour les taux."
            )
        },
        {
            "Element": "Onglet 2 - Graphique Classement",
            "Commentaire": (
                f"Classement horizontal par complétude globale. L'unité la plus performante observée est '{top_comp}'. "
                "Lire du bas vers le haut pour la progression."
            )
        },
        {
            "Element": "Onglet 3 - Tableau Promptitude",
            "Commentaire": (
                f"Promptitude moyenne {prom_avg:.2f}%. Les couleurs suivent la même échelle que l'onglet 3 "
                "pour identifier rapidement les retards de notification."
            )
        },
        {
            "Element": "Onglet 4 - Quadrant Comparatif",
            "Commentaire": (
                "Le quadrant croise complétude (axe X) et promptitude (axe Y). "
                "Les lignes à 80% séparent les unités à renforcer des unités performantes."
            )
        },
        {
            "Element": "Onglet 4 - Tableau comparatif",
            "Commentaire": (
                "Le tableau comparatif présente, par unité, la complétude et la promptitude "
                "avec les compteurs de datasets >=95% pour prioriser les actions."
            )
        },
        {
            "Element": "Onglet 4 - Top/Flop",
            "Commentaire": (
                f"Top complétude: {top_comp}. Flop promptitude: {flop_prom}. "
                "Action recommandée: cibler d'abord les unités du flop promptitude."
            )
        },
        {
            "Element": "Top 5 Complétude",
            "Commentaire": "Classement des 5 unités avec les meilleurs taux de complétude."
        },
        {
            "Element": "Flop 5 Promptitude",
            "Commentaire": "Classement des 5 unités avec les plus faibles taux de promptitude."
        },
        {
            "Element": "Onglet 5 - Tableau Catégorisation",
            "Commentaire": (
                "Le tableau compare M-1 et M: règles violées, règles corrigées et ratio /100 rapports. "
                "Le ratio élevé (texte rouge) signale un risque qualité prioritaire."
            )
        },
    ]

    if df_tab4_fusion is not None and not df_tab4_fusion.empty:
        rows.append(
            {
                "Element": "Onglet 4 - Tableau fusionné (zone filtrée)",
                "Commentaire": (
                    f"Pour '{zone_label}', ce tableau aligne indicateurs dataset Reporting/Actual avec "
                    "complétude, promptitude et compteurs >=95%."
                )
            }
        )
    if df_tab5 is not None and not df_tab5.empty:
        rows.append(
            {
                "Element": "Onglet 5 - Lecture décisionnelle",
                "Commentaire": (
                    "Comparer d'abord 'Règles violées (M)' puis 'Règles corrigées (M-1 -> M)'. "
                    "Un ratio qui baisse avec des règles corrigées en hausse indique une amélioration."
                )
            }
        )

    return pd.DataFrame(rows)

def to_excel_dashboard_report(
    df_base,
    df_actual_detail,
    df_expected_detail,
    df_final_c,
    df_final_p,
    df_synth,
    comments_df,
    df_tab4_fusion=None,
    df_comparatif=None,
    df_tab5=None,
    df_top5=None,
    df_flop5=None,
    style_context=None
):
    """Exporte un rapport Excel commenté avec colorations conditionnelles."""
    style_context = style_context or {}
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        sty_base = _build_styler(df_base, **style_context.get("base", {}))
        (sty_base if sty_base is not None else df_base).to_excel(
            writer, index=False, sheet_name='Base_Donnees'
        )

        sty_actual = _build_styler(df_actual_detail, **style_context.get("actual_detail", {}))
        (sty_actual if sty_actual is not None else df_actual_detail).to_excel(
            writer, index=False, sheet_name='Detail_Donnees_Reelles'
        )

        sty_expected = _build_styler(df_expected_detail, **style_context.get("expected_detail", {}))
        (sty_expected if sty_expected is not None else df_expected_detail).to_excel(
            writer, index=False, sheet_name='Detail_Reports_Attendus'
        )

        sty_comp = _build_styler(df_final_c, **style_context.get("completude", {}))
        (sty_comp if sty_comp is not None else df_final_c).to_excel(
            writer, index=False, sheet_name='Perf_Final_Completude'
        )

        sty_prom = _build_styler(df_final_p, **style_context.get("promptitude", {}))
        (sty_prom if sty_prom is not None else df_final_p).to_excel(
            writer, index=False, sheet_name='Analyse_Promptitude'
        )

        sty_synth = _build_styler(df_synth, **style_context.get("synthese", {}))
        (sty_synth if sty_synth is not None else df_synth).to_excel(
            writer, index=False, sheet_name='Synthese_Comparative'
        )

        if df_comparatif is not None and not df_comparatif.empty:
            sty_compa = _build_styler(df_comparatif, **style_context.get("comparatif", {}))
            (sty_compa if sty_compa is not None else df_comparatif).to_excel(
                writer, index=False, sheet_name='Comparatif_C_P'
            )

        if df_tab4_fusion is not None and not df_tab4_fusion.empty:
            sty_fusion = _build_styler(df_tab4_fusion, **style_context.get("fusion", {}))
            (sty_fusion if sty_fusion is not None else df_tab4_fusion).to_excel(
                writer, index=False, sheet_name='Perf_Zone_Filtree'
            )

        if df_top5 is not None and not df_top5.empty:
            sty_top5 = _build_styler(df_top5, **style_context.get("top5", {}))
            (sty_top5 if sty_top5 is not None else df_top5).to_excel(
                writer, index=False, sheet_name='Top5_Completude'
            )

        if df_flop5 is not None and not df_flop5.empty:
            sty_flop5 = _build_styler(df_flop5, **style_context.get("flop5", {}))
            (sty_flop5 if sty_flop5 is not None else df_flop5).to_excel(
                writer, index=False, sheet_name='Flop5_Promptitude'
            )

        if df_tab5 is not None and not df_tab5.empty:
            sty_tab5 = _build_styler(df_tab5, **style_context.get("tab5", {}))
            (sty_tab5 if sty_tab5 is not None else df_tab5).to_excel(
                writer, index=False, sheet_name='Resultats_Regles'
            )
        else:
            pd.DataFrame(
                [{"Information": "Aucun résultat de règles. Lancez d'abord l'analyse des violations (onglet 5)."}]
            ).to_excel(writer, index=False, sheet_name='Resultats_Regles')

        comments_df.to_excel(writer, index=False, sheet_name='Commentaires')
    return output.getvalue()

def _add_dataframe_to_docx(doc, title, df, max_rows=20):
    doc.add_heading(title, level=2)
    if df is None or df.empty:
        doc.add_paragraph("Aucune donnée disponible.")
        return
    df_show = df.head(max_rows).copy()
    table = doc.add_table(rows=1, cols=len(df_show.columns))
    table.style = 'Table Grid'
    for i, col in enumerate(df_show.columns):
        table.rows[0].cells[i].text = str(col)
    for _, row in df_show.iterrows():
        cells = table.add_row().cells
        for i, col in enumerate(df_show.columns):
            val = row[col]
            if isinstance(val, float):
                cells[i].text = f"{val:.2f}"
            else:
                cells[i].text = str(val)

def _add_dataframe_image_to_docx(doc, title, df, style_ctx=None, comment=None, max_rows=20):
    """Ajoute un tableau en image (avec couleurs) dans Word."""
    from docx.shared import Inches

    doc.add_heading(title, level=2)
    if comment:
        doc.add_paragraph(comment)

    img_bytes, img_error = dataframe_to_png_bytes(
        df,
        title=title,
        max_rows=max_rows,
        style_context=style_ctx
    )
    if img_bytes is not None:
        doc.add_picture(io.BytesIO(img_bytes), width=Inches(6.7))
    else:
        doc.add_paragraph(f"Image indisponible: {img_error}")
        _add_dataframe_to_docx(doc, title, df, max_rows=max_rows)

def to_word_dashboard_report(
    df_final_c,
    df_final_p,
    df_synth,
    comments_df,
    df_tab4_fusion=None,
    fig_comp=None,
    fig_quad=None,
    df_comparatif=None,
    df_tab5=None,
    style_context=None
):
    """Exporte un rapport Word commenté avec tableaux colorés et explications."""
    try:
        from docx import Document
        from docx.shared import Inches
    except Exception:
        return None, "python-docx non installé (pip install python-docx)."

    style_context = style_context or {}
    doc = Document()
    doc.add_heading("Rapport Dashboard SNIS RDC", level=1)
    doc.add_paragraph(f"Généré le {datetime.now().strftime('%Y-%m-%d %H:%M')}")

    doc.add_heading("Commentaires d'interprétation", level=2)
    for _, row in comments_df.iterrows():
        doc.add_paragraph(f"{row['Element']}: {row['Commentaire']}")

    img_comp, _ = figure_to_png_bytes(fig_comp)
    if img_comp is not None:
        doc.add_heading("Graphique - Classement Complétude", level=2)
        doc.add_paragraph(_comment_for(comments_df, "Graphique Classement", "Classement par complétude globale."))
        doc.add_picture(io.BytesIO(img_comp), width=Inches(6.7))

    img_quad, _ = figure_to_png_bytes(fig_quad)
    if img_quad is not None:
        doc.add_heading("Graphique - Quadrant Comparatif", level=2)
        doc.add_paragraph(_comment_for(comments_df, "Quadrant", "Comparaison complétude/promptitude."))
        doc.add_picture(io.BytesIO(img_quad), width=Inches(6.7))

    _add_dataframe_image_to_docx(
        doc,
        "Tableau - Complétude",
        df_final_c,
        style_ctx=style_context.get("completude", {}),
        comment=_comment_for(comments_df, "Tableau Complétude", "")
    )
    _add_dataframe_image_to_docx(
        doc,
        "Tableau - Promptitude",
        df_final_p,
        style_ctx=style_context.get("promptitude", {}),
        comment=_comment_for(comments_df, "Tableau Promptitude", "")
    )
    _add_dataframe_image_to_docx(
        doc,
        "Tableau - Comparatif",
        df_synth,
        style_ctx=style_context.get("synthese", {}),
        comment=_comment_for(comments_df, "Top/Flop", "")
    )
    if df_comparatif is not None and not df_comparatif.empty:
        _add_dataframe_image_to_docx(
            doc,
            "Tableau comparatif Complétude / Promptitude",
            df_comparatif,
            style_ctx=style_context.get("comparatif", {}),
            comment=_comment_for(comments_df, "Tableau comparatif", "")
        )
    if df_tab4_fusion is not None and not df_tab4_fusion.empty:
        _add_dataframe_image_to_docx(
            doc,
            "Tableau fusionné zone filtrée",
            df_tab4_fusion,
            style_ctx=style_context.get("fusion", {}),
            comment=_comment_for(comments_df, "fusionné", "")
        )
    if df_tab5 is not None and not df_tab5.empty:
        _add_dataframe_image_to_docx(
            doc,
            "Tableau catégorisation (M-1 vs M)",
            df_tab5,
            style_ctx=style_context.get("tab5", {}),
            comment=_comment_for(comments_df, "Catégorisation", "")
        )

    output = io.BytesIO()
    doc.save(output)
    return output.getvalue(), None

def dataframe_to_png_bytes(df, title="Tableau", max_rows=20, style_context=None):
    """Convertit un DataFrame en image PNG avec colorations conditionnelles."""
    try:
        import matplotlib.pyplot as plt
    except Exception:
        return None, "matplotlib non installé (pip install matplotlib)."

    if df is None or df.empty:
        return None, "Aucune donnée à convertir en image."

    style_context = style_context or {}
    taux_cols = [c for c in style_context.get("taux_cols", []) if c in df.columns]
    score_cols = [c for c in style_context.get("score_cols", []) if c in df.columns]
    int_cols = [c for c in style_context.get("int_cols", []) if c in df.columns]
    highlight_max_cols = [c for c in style_context.get("highlight_max_cols", []) if c in df.columns]
    ratio_alert_col = style_context.get("ratio_alert_col")
    ratio_alert_threshold = float(style_context.get("ratio_alert_threshold", 10.0))
    decimals = int(style_context.get("decimals", 2))

    df_show = df.head(max_rows).copy()
    max_values = {}
    for c in highlight_max_cols:
        s = pd.to_numeric(df_show[c], errors='coerce')
        max_values[c] = s.max() if not s.isna().all() else np.nan

    cell_text = []
    for _, row in df_show.iterrows():
        row_text = []
        for col in df_show.columns:
            val = row[col]
            if pd.isna(val):
                row_text.append("")
            elif col in int_cols:
                row_text.append(f"{pd.to_numeric(val, errors='coerce'):.0f}")
            elif isinstance(val, (float, np.floating, int, np.integer)):
                row_text.append(f"{float(val):.{decimals}f}")
            else:
                row_text.append(str(val))
        cell_text.append(row_text)

    fig_h = max(2.8, 0.45 * (len(df_show) + 2))
    fig_w = max(10, 1.5 * len(df_show.columns))
    fig, ax = plt.subplots(figsize=(fig_w, fig_h))
    ax.axis('off')
    ax.set_title(title, fontsize=11, pad=8)

    table = ax.table(
        cellText=cell_text,
        colLabels=df_show.columns,
        cellLoc='center',
        loc='center'
    )
    table.auto_set_font_size(False)
    table.set_fontsize(8)
    table.scale(1, 1.2)

    # Style entête
    for j in range(len(df_show.columns)):
        header_cell = table[(0, j)]
        header_cell.set_facecolor("#1F4E78")
        header_cell.set_edgecolor("#FFFFFF")
        header_cell.get_text().set_color("white")
        header_cell.get_text().set_weight("bold")

    # Style cellules data
    for i in range(len(df_show)):
        for j, col in enumerate(df_show.columns):
            val = df_show.iloc[i, j]
            cell = table[(i + 1, j)]
            cell.set_edgecolor("#D9D9D9")
            text_color = "#000000"
            bg_color = "#FFFFFF"
            bold = False

            if col in taux_cols:
                text_color, bg_color, bold_taux = _extract_css_values(style_taux(val))
                bold = bold or bold_taux
            if col in score_cols:
                text_color, bg_color, bold_score = _extract_css_values(style_score(val))
                bold = bold or bold_score

            if col in highlight_max_cols and pd.notna(max_values.get(col)):
                val_num = pd.to_numeric(val, errors='coerce')
                if pd.notna(val_num) and np.isclose(float(val_num), float(max_values[col])):
                    bg_color = "#ffcccc"

            if ratio_alert_col and col == ratio_alert_col:
                val_num = pd.to_numeric(val, errors='coerce')
                if pd.notna(val_num) and float(val_num) > ratio_alert_threshold:
                    text_color = "red"
                    bold = True

            cell.set_facecolor(bg_color)
            cell.get_text().set_color(text_color)
            if bold:
                cell.get_text().set_weight("bold")

    fig.tight_layout()
    output = io.BytesIO()
    fig.savefig(output, format='png', dpi=180, bbox_inches='tight')
    plt.close(fig)
    return output.getvalue(), None

def figure_to_png_bytes(fig):
    """Convertit un graphique Plotly en PNG (retourne bytes, error)."""
    if fig is None:
        return None, "Figure indisponible."
    try:
        return fig.to_image(format="png", scale=2), None
    except Exception:
        return None, "Export image Plotly indisponible (installer kaleido: pip install kaleido)."

def _comment_for(comments_df, contains_text, default_text):
    if comments_df is None or comments_df.empty:
        return default_text
    mask = comments_df["Element"].astype(str).str.contains(contains_text, case=False, na=False)
    if mask.any():
        return str(comments_df.loc[mask, "Commentaire"].iloc[0])
    return default_text

def _add_image_slide(prs, title, comment, image_bytes):
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    slide.shapes.title.text = title

    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12.3), Inches(1.2))
    tx_frame = tx_box.text_frame
    tx_frame.text = comment
    tx_frame.word_wrap = True
    tx_frame.paragraphs[0].font.size = Pt(12)

    if image_bytes is not None:
        slide.shapes.add_picture(io.BytesIO(image_bytes), Inches(0.5), Inches(1.9), width=Inches(12.3))

def to_powerpoint_dashboard_report(
    df_synth,
    comments_df,
    fig_comp=None,
    fig_quad=None,
    df_comparatif=None,
    df_tab4_fusion=None,
    df_final_c=None,
    df_final_p=None,
    df_tab5=None,
    df_base=None,
    df_actual_detail=None,
    df_expected_detail=None,
    df_top5=None,
    df_flop5=None,
    style_context=None
):
    """Exporte un PowerPoint avec graphiques + tableaux colorés + commentaires clairs."""
    try:
        from pptx import Presentation
    except Exception:
        return None, "python-pptx non installé (pip install python-pptx)."

    style_context = style_context or {}
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Rapport Dashboard SNIS RDC"
    slide.placeholders[1].text = f"Généré le {datetime.now().strftime('%Y-%m-%d %H:%M')}"

    # 1) Graphique complétude
    img_comp, _ = figure_to_png_bytes(fig_comp)
    _add_image_slide(
        prs,
        "Graphique Classement Complétude",
        _comment_for(comments_df, "Graphique Classement", "Classement des unités par complétude globale."),
        img_comp
    )

    # 1-bis) Tableau complétude coloré
    if df_final_c is not None and not df_final_c.empty:
        img_tab_comp, _ = dataframe_to_png_bytes(
            df_final_c,
            title="Tableau Complétude",
            style_context=style_context.get("completude", {})
        )
        _add_image_slide(
            prs,
            "Tableau Complétude (coloré)",
            _comment_for(comments_df, "Tableau Complétude", "Lecture détaillée de la complétude par unité."),
            img_tab_comp
        )

    # 1-ter) Base de données
    if df_base is not None and not df_base.empty:
        img_base, _ = dataframe_to_png_bytes(
            df_base,
            title="Base de données",
            style_context=style_context.get("base", {})
        )
        _add_image_slide(
            prs,
            "Base de données",
            _comment_for(comments_df, "Base de données", "Données sources utilisées après filtres."),
            img_base
        )

    # 1-quater) Rapports détaillés données réelles
    if df_actual_detail is not None and not df_actual_detail.empty:
        img_actual, _ = dataframe_to_png_bytes(
            df_actual_detail,
            title="Rapports détaillés - Données réelles",
            style_context=style_context.get("actual_detail", {})
        )
        _add_image_slide(
            prs,
            "Rapports détaillés - Données réelles",
            _comment_for(comments_df, "Données réelles", "Détail des actual reports et total Reports_Actual."),
            img_actual
        )

    # 1-quint) Rapports détaillés participants/attendus
    if df_expected_detail is not None and not df_expected_detail.empty:
        img_expected, _ = dataframe_to_png_bytes(
            df_expected_detail,
            title="Rapports détaillés - Participants/Attendus",
            style_context=style_context.get("expected_detail", {})
        )
        _add_image_slide(
            prs,
            "Rapports détaillés - Participants/Attendus",
            _comment_for(comments_df, "Participants/Attendus", "Détail des expected reports et total Reports_Attendu."),
            img_expected
        )

    # 2) Quadrant comparatif
    img_quad, _ = figure_to_png_bytes(fig_quad)
    _add_image_slide(
        prs,
        "Quadrant Complétude / Promptitude",
        _comment_for(comments_df, "Quadrant", "Comparaison de la complétude et de la promptitude."),
        img_quad
    )

    # 2-bis) Tableau promptitude coloré
    if df_final_p is not None and not df_final_p.empty:
        img_tab_prom, _ = dataframe_to_png_bytes(
            df_final_p,
            title="Tableau Promptitude",
            style_context=style_context.get("promptitude", {})
        )
        _add_image_slide(
            prs,
            "Tableau Promptitude (coloré)",
            _comment_for(comments_df, "Tableau Promptitude", "Lecture détaillée de la promptitude par unité."),
            img_tab_prom
        )

    # 3) Tableau comparatif complétude/promptitude
    img_comp_table, _ = dataframe_to_png_bytes(
        df_comparatif,
        title="Tableau comparatif Complétude / Promptitude",
        style_context=style_context.get("comparatif", {})
    )
    _add_image_slide(
        prs,
        "Tableau comparatif",
        _comment_for(comments_df, "Tableau comparatif", "Lecture comparative par unité."),
        img_comp_table
    )

    # 4) Tableau fusionné zone filtrée (si disponible)
    if df_tab4_fusion is not None and not df_tab4_fusion.empty:
        img_fusion, _ = dataframe_to_png_bytes(
            df_tab4_fusion,
            title="Tableau fusionné indicateurs dataset",
            style_context=style_context.get("fusion", {})
        )
        _add_image_slide(
            prs,
            "Tableau fusionné zone filtrée",
            _comment_for(comments_df, "fusionné", "Vue détaillée des indicateurs dataset en zone filtrée."),
            img_fusion
        )

    # 4-bis) Top 5 complétude
    if df_top5 is not None and not df_top5.empty:
        img_top5, _ = dataframe_to_png_bytes(
            df_top5,
            title="Top 5 Complétude",
            style_context=style_context.get("top5", {})
        )
        _add_image_slide(
            prs,
            "Top 5 Complétude",
            _comment_for(comments_df, "Top 5 Complétude", "Unités les plus performantes en complétude."),
            img_top5
        )

    # 4-ter) Flop 5 promptitude
    if df_flop5 is not None and not df_flop5.empty:
        img_flop5, _ = dataframe_to_png_bytes(
            df_flop5,
            title="Flop 5 Promptitude",
            style_context=style_context.get("flop5", {})
        )
        _add_image_slide(
            prs,
            "Flop 5 Promptitude",
            _comment_for(comments_df, "Flop 5 Promptitude", "Unités prioritaires à renforcer en promptitude."),
            img_flop5
        )

    # 5) Tableau catégorisation (si disponible)
    if df_tab5 is not None and not df_tab5.empty:
        img_tab5, _ = dataframe_to_png_bytes(
            df_tab5,
            title="Tableau Catégorisation (M-1 vs M)",
            style_context=style_context.get("tab5", {})
        )
        _add_image_slide(
            prs,
            "Tableau Catégorisation (coloré)",
            _comment_for(comments_df, "Catégorisation", "Lecture des violations, corrections et ratio /100 rapports."),
            img_tab5
        )
    else:
        _add_image_slide(
            prs,
            "Résultats des règles de validation",
            "Aucun résultat à afficher. Lancez d'abord l'analyse des violations dans l'onglet 5.",
            None
        )

    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Synthèse"
    comp_avg = float(df_synth['Complétude_Globale (%)'].mean()) if not df_synth.empty else 0.0
    prom_avg = float(df_synth['Promptitude_Globale (%)'].mean()) if not df_synth.empty else 0.0
    s.placeholders[1].text = (
        f"Complétude moyenne: {comp_avg:.2f}%\n"
        f"Promptitude moyenne: {prom_avg:.2f}%\n"
        f"Unités analysées: {len(df_synth)}"
    )

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue(), None

# Configuration pour téléchargement d'images Plotly
plotly_config = {
    'toImageButtonOptions': {
        'format': 'png',
        'filename': 'export_graphique_snis',
        'scale': 2
    }
}

# --- 2. FONCTIONS API DHIS2 ---
def _read_config_value(key):
    """Lit une valeur de config depuis Streamlit secrets, sinon variables d'environnement."""
    value = None
    try:
        value = st.secrets.get(key)
    except Exception:
        value = None
    if value is None or str(value).strip() == "":
        value = os.getenv(key)
    return str(value).strip() if value is not None else None

def _read_int_config_value(key, default_value):
    """Lit un entier de config; retourne default_value si absent/invalide."""
    raw_value = _read_config_value(key)
    if raw_value is None or str(raw_value).strip() == "":
        return default_value
    try:
        parsed = int(str(raw_value).strip())
        return parsed if parsed > 0 else default_value
    except Exception:
        return default_value

def test_dhis2_credentials(base_url, username, password):
    """Teste les identifiants utilisateur DHIS2 sur /api/me."""
    try:
        response = requests.get(
            f"{base_url.rstrip('/')}/api/me",
            auth=(username, password),
            timeout=10
        )
    except Exception as e:
        return False, f"Erreur de connexion : {e}"

    if response.status_code == 200:
        return True, None
    if response.status_code in (401, 403):
        return False, "Identifiants incorrects. Veuillez réessayer."
    return False, f"Connexion DHIS2 échouée (HTTP {response.status_code})."

@st.cache_resource
def _build_dhis2_client(url, username, password, timeout_connect=10, timeout_read=90, retries=2, backoff=1.0):
    class SimpleDHIS2Client:
        """Client DHIS2 minimal sans dépendance openhexa."""
        def __init__(self, base_url, user, pwd, timeout_connect_sec=10, timeout_read_sec=90, retries_count=2, backoff_factor=1.0):
            self.base_url = str(base_url).rstrip("/")
            self.timeout = (timeout_connect_sec, timeout_read_sec)
            self.session = requests.Session()
            self.session.auth = (user, pwd)
            self.session.headers.update({"Accept": "application/json"})

            retry_strategy = Retry(
                total=retries_count,
                connect=retries_count,
                read=retries_count,
                status=retries_count,
                backoff_factor=backoff_factor,
                status_forcelist=[429, 500, 502, 503, 504],
                allowed_methods=["GET", "HEAD", "OPTIONS"],
                raise_on_status=False,
            )
            adapter = HTTPAdapter(max_retries=retry_strategy)
            self.session.mount("https://", adapter)
            self.session.mount("http://", adapter)

        def get(self, endpoint, params=None):
            clean_endpoint = str(endpoint).lstrip("/")
            if not clean_endpoint.startswith("api/"):
                clean_endpoint = f"api/{clean_endpoint}"
            url = f"{self.base_url}/{clean_endpoint}"
            response = self.session.get(url, params=params, timeout=self.timeout)
            if response.status_code in (401, 403):
                raise RuntimeError("Accès DHIS2 refusé (vérifie les identifiants et permissions).")
            response.raise_for_status()
            if not response.text:
                return {}
            return response.json()

    return SimpleDHIS2Client(url, username, password, timeout_connect, timeout_read, retries, backoff)

def get_dhis2_client():
    dhis2_url = _read_config_value("DHIS2_URL")
    dhis2_user = str(st.session_state.get("dhis2_user", "")).strip()
    dhis2_pass = str(st.session_state.get("dhis2_pass", "")).strip()
    dhis2_timeout_connect = _read_int_config_value("DHIS2_TIMEOUT_CONNECT", 10)
    dhis2_timeout_read = _read_int_config_value("DHIS2_TIMEOUT_READ", 90)
    dhis2_retries = _read_int_config_value("DHIS2_HTTP_RETRIES", 2)

    missing_keys = []
    if not dhis2_url:
        missing_keys.append("DHIS2_URL")
    if not dhis2_user:
        missing_keys.append("Nom d'utilisateur DHIS2")
    if not dhis2_pass:
        missing_keys.append("Mot de passe DHIS2")
    if missing_keys:
        missing_text = ", ".join(missing_keys)
        raise RuntimeError(
            f"Configuration DHIS2 manquante: {missing_text}. "
            "Ajoute DHIS2_URL dans Streamlit Cloud (App settings > Secrets), "
            "puis connecte-toi via la barre latérale."
        )

    return _build_dhis2_client(
        dhis2_url,
        dhis2_user,
        dhis2_pass,
        timeout_connect=dhis2_timeout_connect,
        timeout_read=dhis2_timeout_read,
        retries=dhis2_retries,
        backoff=1.0,
    )

@st.cache_data(ttl=3600)
def get_data(favori_id, period=None, cache_user=None):
    _ = cache_user
    if not favori_id: return None
    try:
        dhis = get_dhis2_client()
        endpoint = f"visualizations/{favori_id}/data.json"
        if period:
            endpoint += f"?dimension=pe:{period}"

        response = dhis.get(endpoint)
        headers, rows = response.get("headers", []), response.get("rows", [])
        columns = [h.get("name", h.get("column")) for h in headers]
        return pd.DataFrame(rows, columns=columns)
    except requests.exceptions.ReadTimeout:
        st.error(
            "Erreur de connexion : délai de lecture dépassé sur DHIS2. "
            "Essaie une période plus courte (ex. 1-2 mois) ou augmente "
            "DHIS2_TIMEOUT_READ dans les secrets Streamlit."
        )
        return None
    except requests.exceptions.Timeout:
        st.error(
            "Erreur de connexion : timeout réseau DHIS2. "
            "Vérifie la connectivité et réessaie."
        )
        return None
    except Exception as e:
        st.error(f"Erreur de connexion : {e}")
        return None

@st.cache_data(ttl=3600)
def get_validation_groups(cache_user=None):
    _ = cache_user
    try:
        dhis = get_dhis2_client()
        return dhis.get("validationRuleGroups", params={"fields": "id,displayName", "paging": "false"})['validationRuleGroups']
    except: return []

@st.cache_data(ttl=3600)
def get_validation_results(ou_id, period_list, group_id, cache_user=None):
    _ = cache_user
    # Sécurité anti-Erreur 409 : si l'ID n'est pas un ID système DHIS2, on bloque.
    if not ou_id or ou_id == "USER_ORGUNIT" or len(str(ou_id)) < 5:
        return pd.DataFrame()

    all_results = []
    try:
        dhis = get_dhis2_client()
    except Exception as e:
        st.error(f"Erreur de connexion DHIS2 : {e}")
        return pd.DataFrame()
    for pe in period_list:
        try:
            params = {"ou": ou_id, "pe": pe, "ouMode": "DESCENDANTS"}
            if group_id:
                params["vrg"] = group_id
            res = dhis.get("validationResults", params=params)
            if 'validationResults' in res:
                all_results.extend(res['validationResults'])
        except Exception:
            continue # Ignore les erreurs de périodes vides
    return pd.DataFrame(all_results)

@st.cache_data(ttl=3600)
def get_children_org_units_details(parent_id, cache_user=None):
    _ = cache_user
    """Récupère les enfants directs (id + nom) d'une unité d'organisation."""
    if not parent_id:
        return []
    try:
        dhis = get_dhis2_client()
        response = dhis.get(
            f"organisationUnits/{parent_id}",
            params={"fields": "children[id,displayName]", "paging": "false"}
        )
        return response.get("children", []) or []
    except Exception:
        return []

@st.cache_data(ttl=3600)
def get_org_units_hierarchy(org_unit_ids, cache_user=None):
    """
    Récupère la hiérarchie organisationnelle (pays/province/zone) pour une liste d'IDs.
    Retour: dict[ou_id] -> {
        country_name, country_id, province_name, province_id, zone_name, zone_id, level
    }
    """
    _ = cache_user
    cleaned_ids = []
    for ou_id in (org_unit_ids or []):
        if ou_id is None:
            continue
        ou_value = str(ou_id).strip()
        if not ou_value or ou_value.lower() in {"nan", "none", "null"}:
            continue
        cleaned_ids.append(ou_value)
    unique_ids = sorted(set(cleaned_ids))
    if not unique_ids:
        return {}

    try:
        dhis = get_dhis2_client()
    except Exception:
        return {}

    hierarchy = {}
    chunk_size = 60
    for i in range(0, len(unique_ids), chunk_size):
        chunk = unique_ids[i:i + chunk_size]
        try:
            params = {
                "fields": "id,displayName,level,ancestors[id,displayName,level]",
                "filter": f"id:in:[{','.join(chunk)}]",
                "paging": "false",
            }
            response = dhis.get("organisationUnits", params=params)
            units = response.get("organisationUnits", [])
        except Exception:
            units = []

        for ou in units:
            ou_id = str(ou.get("id", "")).strip()
            if not ou_id:
                continue
            ou_name = str(ou.get("displayName", "")).strip()
            ou_level = ou.get("level")

            ancestors = ou.get("ancestors", []) or []
            ancestors_sorted = sorted(
                [a for a in ancestors if isinstance(a, dict) and a.get("id")],
                key=lambda a: int(a.get("level") or 0)
            )

            chain = []
            seen = set()
            for item in ancestors_sorted + [{"id": ou_id, "displayName": ou_name, "level": ou_level}]:
                item_id = str(item.get("id", "")).strip()
                if not item_id or item_id in seen:
                    continue
                seen.add(item_id)
                chain.append(item)

            if not chain:
                continue

            country = chain[0]
            province = chain[1] if len(chain) > 1 else chain[0]
            zone = chain[2] if len(chain) > 2 else (chain[1] if len(chain) > 1 else chain[0])

            hierarchy[ou_id] = {
                "country_name": str(country.get("displayName", "")).strip(),
                "country_id": str(country.get("id", "")).strip(),
                "province_name": str(province.get("displayName", "")).strip(),
                "province_id": str(province.get("id", "")).strip(),
                "zone_name": str(zone.get("displayName", "")).strip(),
                "zone_id": str(zone.get("id", "")).strip(),
                "level": int(ou_level) if ou_level is not None else None,
            }

    return hierarchy

# --- 3. SIDEBAR & FILTRES ---
with st.sidebar:
    st.image("https://snisrdc.com/dhis-web-commons/security/logo_front.png", width=150)
    st.title("⚙️ Configuration")

    st.subheader("🔐 Connexion DHIS2")
    base_url = _read_config_value("DHIS2_URL")
    if not base_url:
        st.error("DHIS2_URL manquant. Ajoute-le dans les secrets Streamlit Cloud.")
        st.stop()

    user_input = st.text_input(
        "Nom d'utilisateur",
        value=str(st.session_state.get("dhis2_user", "")),
        key="dhis2_user_input"
    )
    password_input = st.text_input(
        "Mot de passe",
        type="password",
        key="dhis2_password_input"
    )

    col_conn, col_disc = st.columns(2)
    login_button = col_conn.button("Se connecter", use_container_width=True)
    logout_button = col_disc.button("Se déconnecter", use_container_width=True)

    if logout_button:
        for key in ["connected", "dhis2_user", "dhis2_pass", "dhis2_user_input", "dhis2_password_input"]:
            st.session_state.pop(key, None)
        st.cache_resource.clear()
        st.rerun()

    if login_button:
        if not user_input or not password_input:
            st.session_state["connected"] = False
            st.warning("Renseigne le nom d'utilisateur et le mot de passe.")
        else:
            ok, error_msg = test_dhis2_credentials(base_url, user_input, password_input)
            if ok:
                st.cache_resource.clear()
                st.session_state["connected"] = True
                st.session_state["dhis2_user"] = user_input
                st.session_state["dhis2_pass"] = password_input
                st.success(f"Connecté en tant que : {user_input}")
                st.rerun()
            else:
                st.session_state["connected"] = False
                st.error(error_msg)

    if st.session_state.get("connected"):
        st.success(f"Connecté en tant que : {st.session_state.get('dhis2_user')}")
    else:
        st.info("Veuillez vous connecter via la barre latérale pour accéder aux données.")
        st.stop()

    st.divider()

    dict_favoris = {
        "Performance Globale (ROzCY14OLTE)": "ROzCY14OLTE",
        "Analyse Promptitude (mldsgxAvIIi)": "mldsgxAvIIi",
        "Rapport TyAJb0qitMz": "TyAJb0qitMz",
        "Autre ID personnalisé": "CUSTOM"
    }
    choix_fav = st.selectbox("Rapport DHIS2 :", list(dict_favoris.keys()))
    id_final = st.text_input("Saisissez l'ID :") if choix_fav == "Autre ID personnalisé" else dict_favoris[choix_fav]

    st.divider()
    st.subheader("📅 Sélection de la Période")
    annee_choisie = st.selectbox("Choisir l'année :", ["2026", "2025", "2024", "2023"], index=1)

    mois_liste = ["Janvier", "Février", "Mars", "Avril", "Mai", "Juin", "Juillet", "Août", "Septembre", "Octobre", "Novembre", "Décembre"]
    mois_dict = {m: str(i+1).zfill(2) for i, m in enumerate(mois_liste)}

    col_m1, col_m2 = st.columns(2)
    with col_m1: mois_debut = st.selectbox("Mois début :", mois_liste, index=0)
    with col_m2: mois_fin = st.selectbox("Mois fin :", mois_liste, index=0)

    idx_start, idx_end = mois_liste.index(mois_debut), mois_liste.index(mois_fin)

    # Valeurs par défaut pour éviter des variables non définies.
    selection_mois = [f"{annee_choisie}{mois_dict[mois_debut]}"]
    period_id = selection_mois[0]
    prev_idx = idx_start - 1 if idx_start > 0 else 11
    prev_year = annee_choisie if idx_start > 0 else str(int(annee_choisie)-1)
    selection_mois_prev = [f"{prev_year}{mois_dict[mois_liste[prev_idx]]}"]

    if idx_start <= idx_end:
        selection_mois = [f"{annee_choisie}{mois_dict[mois_liste[i]]}" for i in range(idx_start, idx_end + 1)]
        period_id = ";".join(selection_mois)
    else:
        st.error("Le mois de début doit être avant le mois de fin.")

# --- 4. CHARGEMENT & FILTRAGE DES DONNÉES ---
current_user_for_cache = st.session_state.get("dhis2_user", "")
df_raw = get_data(id_final, period=period_id, cache_user=current_user_for_cache)
if df_raw is not None:
    df_raw = normalize_orgunit_columns(df_raw)

if df_raw is not None:
    # Récupération dynamique de l'ID de l'unité d'organisation (OrgUnit ID)
    # On cherche l'ID dans les colonnes renvoyées par le favori
    mapping_ou_id = {}
    if 'Organisation unit' in df_raw.columns and 'Organisation unit ID' in df_raw.columns:
        mapping_ou_id = dict(zip(df_raw['Organisation unit'], df_raw['Organisation unit ID']))

    # Indicateur parent (zone/province) pour filtrer les aires de santé.
    is_parent_by_ou = {}
    if 'Organisation unit' in df_raw.columns and 'Organisation unit is parent' in df_raw.columns:
        temp_parent = df_raw[['Organisation unit', 'Organisation unit is parent']].drop_duplicates('Organisation unit')
        is_parent_by_ou = {
            row['Organisation unit']: str(row['Organisation unit is parent']).strip().lower() in ['true', '1']
            for _, row in temp_parent.iterrows()
        }

    target_df_all = df_raw.copy()
    default_cols = ['Organisation unit code', 'Organisation unit description',
                    'Reporting month', 'Organisation unit parameter', 'Organisation unit is parent']
    cols_existing = [c for c in default_cols if c in target_df_all.columns]
    if cols_existing: target_df_all.drop(columns=cols_existing, inplace=True)

    st.sidebar.subheader("🔍 Filtrage")
    scope_df = pd.DataFrame({"Organisation unit": df_raw.get("Organisation unit", pd.Series(dtype=str))})
    if "Organisation unit ID" in df_raw.columns:
        scope_df["Organisation unit ID"] = df_raw["Organisation unit ID"].astype(str)
    else:
        scope_df["Organisation unit ID"] = ""

    hierarchy_by_id = get_org_units_hierarchy(
        scope_df["Organisation unit ID"].dropna().tolist(),
        cache_user=current_user_for_cache
    )

    scope_df["Pays"] = scope_df["Organisation unit ID"].map(
        lambda x: hierarchy_by_id.get(str(x), {}).get("country_name", "")
    )
    scope_df["Pays_ID"] = scope_df["Organisation unit ID"].map(
        lambda x: hierarchy_by_id.get(str(x), {}).get("country_id", "")
    )
    scope_df["Province"] = scope_df["Organisation unit ID"].map(
        lambda x: hierarchy_by_id.get(str(x), {}).get("province_name", "")
    )
    scope_df["Province_ID"] = scope_df["Organisation unit ID"].map(
        lambda x: hierarchy_by_id.get(str(x), {}).get("province_id", "")
    )
    scope_df["Zone"] = scope_df["Organisation unit ID"].map(
        lambda x: hierarchy_by_id.get(str(x), {}).get("zone_name", "")
    )
    scope_df["Zone_ID"] = scope_df["Organisation unit ID"].map(
        lambda x: hierarchy_by_id.get(str(x), {}).get("zone_id", "")
    )

    scope_df["Pays"] = scope_df["Pays"].replace("", np.nan).fillna("Pays")
    scope_df["Province"] = scope_df["Province"].replace("", np.nan).fillna(scope_df["Organisation unit"])
    scope_df["Zone"] = scope_df["Zone"].replace("", np.nan).fillna(scope_df["Organisation unit"])

    province_to_id = (
        scope_df[scope_df["Province_ID"].astype(str).str.strip() != ""]
        .drop_duplicates("Province")
        .set_index("Province")["Province_ID"]
        .to_dict()
    )
    zone_to_id = (
        scope_df[scope_df["Zone_ID"].astype(str).str.strip() != ""]
        .drop_duplicates("Zone")
        .set_index("Zone")["Zone_ID"]
        .to_dict()
    )

    scope_df["is_zone_candidate"] = scope_df["Zone"].astype(str) != scope_df["Province"].astype(str)
    if not scope_df["is_zone_candidate"].any() and is_parent_by_ou:
        scope_df["is_zone_candidate"] = ~scope_df["Organisation unit"].map(is_parent_by_ou).fillna(False)
    if not scope_df["is_zone_candidate"].any():
        scope_df["is_zone_candidate"] = True

    filter_level = st.sidebar.selectbox(
        "Niveau d'analyse :",
        ["Synthèse pays", "Province", "Zone de santé"],
        key="filter_level_main"
    )

    selected_province = "Toutes les provinces"
    selected_zone = "Toutes les zones"
    selected_scope_label = "Synthèse pays"
    current_id_systeme = None

    provinces = sorted(scope_df["Province"].dropna().unique().tolist())
    country_names = sorted(scope_df["Pays"].dropna().unique().tolist())
    country_label = country_names[0] if country_names else "Pays"

    if filter_level == "Synthèse pays":
        target_df = target_df_all.copy()
        selected_scope_label = f"Synthèse {country_label}"
        country_ids = [cid for cid in scope_df["Pays_ID"].dropna().astype(str).tolist() if cid.strip()]
        current_id_systeme = country_ids[0] if country_ids else None

    elif filter_level == "Province":
        if not provinces:
            st.warning("Aucune province détectée dans les données.")
            target_df = target_df_all.copy()
            selected_scope_label = "Province indisponible"
        else:
            selected_province = st.sidebar.selectbox("Province :", provinces, key="filter_province_main")
            mask_province = scope_df["Province"] == selected_province
            target_df = target_df_all[mask_province.values].copy()
            current_id_systeme = province_to_id.get(selected_province)
            selected_scope_label = f"Province {selected_province}"

    else:
        province_options = ["Toutes les provinces"] + provinces if provinces else ["Toutes les provinces"]
        selected_province = st.sidebar.selectbox("Province :", province_options, key="filter_zone_province_main")

        zone_pool = scope_df.copy()
        if selected_province != "Toutes les provinces":
            zone_pool = zone_pool[zone_pool["Province"] == selected_province]

        zone_options = sorted(
            zone_pool[zone_pool["is_zone_candidate"]]["Zone"].dropna().unique().tolist()
        )
        if not zone_options:
            zone_options = sorted(zone_pool["Organisation unit"].dropna().unique().tolist())

        if not zone_options:
            st.warning("Aucune zone de santé détectée pour ce filtre.")
            target_df = target_df_all.iloc[0:0].copy()
            selected_scope_label = "Zone indisponible"
        else:
            selected_zone = st.sidebar.selectbox("Zone de santé :", zone_options, key="filter_zone_main")
            mask_exact = target_df_all["Organisation unit"] == selected_zone
            if mask_exact.any():
                target_df = target_df_all[mask_exact].copy()
            else:
                target_df = target_df_all[(scope_df["Zone"] == selected_zone).values].copy()
            current_id_systeme = zone_to_id.get(selected_zone) or mapping_ou_id.get(selected_zone)
            selected_scope_label = f"Zone {selected_zone}"

    # Garde-fou final: garantir la colonne 'Organisation unit' après tous les filtres.
    if "Organisation unit" not in target_df_all.columns:
        target_df_all = normalize_orgunit_columns(target_df_all)
    if "Organisation unit" not in target_df.columns:
        target_df = normalize_orgunit_columns(target_df)

    if "Organisation unit" not in target_df.columns:
        fallback_org_col = _find_best_column(
            target_df.columns.tolist(),
            include_terms=["organisation", "unit"]
        ) or _find_best_column(
            target_df.columns.tolist(),
            include_terms=["org", "unit"]
        )

        if fallback_org_col and fallback_org_col in target_df.columns:
            target_df = target_df.rename(columns={fallback_org_col: "Organisation unit"})
        elif target_df.shape[1] > 0:
            # Dernier recours: créer la colonne à partir de la première colonne disponible.
            target_df = target_df.copy()
            target_df["Organisation unit"] = target_df.iloc[:, 0].astype(str)
            st.warning("La colonne 'Organisation unit' n'était pas fournie par le favori. Un fallback a été appliqué.")
        else:
            st.error("Impossible d'identifier les unités d'organisation dans ce favori DHIS2.")
            st.stop()

    if target_df.empty:
        st.warning("Aucune donnée trouvée avec ce niveau de filtre.")

    is_zone_focus = filter_level == "Zone de santé" and selected_zone != "Toutes les zones"

    st.title(f"📊 SNIS RDC - Performance {selected_scope_label}")

    # Variables partagées pour export hors onglets.
    fig_comp = None
    fig_quad = None
    df_tab4_fusion = pd.DataFrame()
    df_comparatif = pd.DataFrame()
    df_tab5_export = pd.DataFrame()
    df_tab1_export = pd.DataFrame()
    df_actual_detail_export = pd.DataFrame()
    df_expected_detail_export = pd.DataFrame()
    df_top5_export = pd.DataFrame()
    df_flop5_export = pd.DataFrame()

    tab1, tab2, tab3, tab4, tab5 = st.tabs(["📁 Base de données", "✅ Complétude", "⏱️ Promptitude", "⚖️ Analyse Comparative", "🩺 Éléments de catégorisation"])

    # --- ONGLET 1 : BASE DE DONNÉES ---
    with tab1:
        st.header("Présentation des données brutes")
        st.write(f"Période affichée : **{mois_debut} à {mois_fin} {annee_choisie}**")
        df_tab1 = target_df.copy()
        if 'Organisation unit ID' in df_tab1.columns:
            df_tab1 = df_tab1.drop(columns=['Organisation unit ID'])
        if is_parent_by_ou and 'Organisation unit' in df_tab1.columns:
            df_tab1 = df_tab1[~df_tab1['Organisation unit'].map(is_parent_by_ou).fillna(False)]

        if df_tab1.empty:
            st.warning("Aucune aire de santé à afficher avec ce filtre.")
        else:
            st.dataframe(df_tab1.head(15), use_container_width=True)
        df_tab1_export = df_tab1.copy()

    # --- ONGLET 2 : COMPLÉTUDE ---
    with tab2:
        st.header("Analyse de la Complétude")
        col_actual_strict = [c for c in target_df.columns if 'actual reports' in c.lower() and 'time' not in c.lower()]
        if col_actual_strict:
            df_actual = target_df[['Organisation unit'] + col_actual_strict].copy()
            for col in col_actual_strict: df_actual[col] = pd.to_numeric(df_actual[col], errors='coerce').fillna(0)
            df_actual['Reports_Actual'] = df_actual[col_actual_strict].sum(axis=1).round(2)
            with st.expander("Détail Reports Actuals"): st.dataframe(df_actual.head(15))
            df_actual_detail_export = df_actual.copy()

        col_expected_strict = [c for c in target_df.columns if 'expected reports' in c.lower() and 'time' not in c.lower()]
        if col_expected_strict:
            df_expected = target_df[['Organisation unit'] + col_expected_strict].copy()
            for col in col_expected_strict: df_expected[col] = pd.to_numeric(df_expected[col], errors='coerce').fillna(0)
            df_expected['Reports_Attendu'] = df_expected[col_expected_strict].sum(axis=1).round(2)
            with st.expander("Détail Reports Attendus"): st.dataframe(df_expected.head(15))
            df_expected_detail_export = df_expected.copy()

        col_rate_uniquement = [c for c in target_df.columns if 'reporting rate' in c.lower() and 'on time' not in c.lower()]
        df_affichage = target_df[['Organisation unit']].copy()
        if col_rate_uniquement:
            df_affichage = target_df[['Organisation unit'] + col_rate_uniquement].copy()
            for col in col_rate_uniquement: df_affichage[col] = pd.to_numeric(df_affichage[col], errors='coerce').fillna(0).round(2)
            with st.expander(f"Affichage de {len(col_rate_uniquement)} indicateurs"): st.dataframe(df_affichage.head(15), use_container_width=True)

        df_synthese = target_df[['Organisation unit']].copy()
        df_synthese['Reports_Actual'] = target_df[col_actual_strict].apply(pd.to_numeric, errors='coerce').fillna(0).sum(axis=1).round(2) if col_actual_strict else 0
        df_synthese['Reports_Attendu'] = target_df[col_expected_strict].apply(pd.to_numeric, errors='coerce').fillna(0).sum(axis=1).round(2) if col_expected_strict else 0
        df_synthese['Complétude_Globale (%)'] = (df_synthese['Reports_Actual'] / df_synthese['Reports_Attendu'].replace(0, np.nan) * 100).fillna(0).round(2)

        st.subheader("Tableau de Performance Final (Stylisé)")
        df_final_c = pd.merge(df_affichage, df_synthese[['Organisation unit', 'Reports_Actual', 'Reports_Attendu', 'Complétude_Globale (%)']], on='Organisation unit')
        df_final_c['Nombre des data set complétude >/=95%'] = (df_final_c[col_rate_uniquement] >= 95).sum(axis=1) if col_rate_uniquement else 0

        st.dataframe(
            df_final_c.style.format({c: "{:.2f}" for c in df_final_c.columns if c not in ['Organisation unit', 'Nombre des data set complétude >/=95%']})
            .map(style_taux, subset=col_rate_uniquement + ['Complétude_Globale (%)'])
            .map(style_score, subset=['Nombre des data set complétude >/=95%']),
            use_container_width=True
        )

        st.divider()
        fig_comp = px.bar(df_final_c.sort_values('Complétude_Globale (%)'), x='Complétude_Globale (%)', y='Organisation unit', orientation='h', color='Complétude_Globale (%)', color_continuous_scale='RdYlGn', title="Classement des zones")
        st.plotly_chart(fig_comp, use_container_width=True, config=plotly_config, key="chart_comp_tab2")

    # --- ONGLET 3 : PROMPTITUDE ---
    with tab3:
        st.header("Analyse de la Promptitude")
        col_rate_ot = [c for c in target_df.columns if 'reporting rate' in c.lower() and 'on time' in c.lower()]
        actual_ot_cols = [c for c in target_df.columns if 'actual reports' in c.lower() and 'on time' in c.lower()]

        df_final_p = target_df[['Organisation unit']].copy()
        for col in col_rate_ot: df_final_p[col] = pd.to_numeric(target_df[col], errors='coerce').fillna(0).round(2)
        df_final_p['Reports_Actual_On_Time'] = target_df[actual_ot_cols].apply(pd.to_numeric, errors='coerce').fillna(0).sum(axis=1).round(2) if actual_ot_cols else 0
        df_final_p['Reports_Attendu'] = target_df[col_expected_strict].apply(pd.to_numeric, errors='coerce').fillna(0).sum(axis=1).round(2) if col_expected_strict else 0
        df_final_p['Promptitude_Globale (%)'] = (df_final_p['Reports_Actual_On_Time'] / df_final_p['Reports_Attendu'].replace(0, np.nan) * 100).fillna(0).round(2)

        nom_col_score_p = 'Nombre des data set promptitude >/=95%'
        df_final_p[nom_col_score_p] = (df_final_p[col_rate_ot] >= 95).sum(axis=1) if col_rate_ot else 0

        st.dataframe(
            df_final_p.style.format({c: "{:.2f}" for c in df_final_p.columns if c not in ['Organisation unit', nom_col_score_p]})
            .map(style_taux, subset=col_rate_ot + ['Promptitude_Globale (%)'])
            .map(style_score, subset=[nom_col_score_p]), use_container_width=True
        )

    # --- ONGLET 4 : ANALYSE COMPARATIVE ---
    with tab4:
        st.header("⚖️ Analyse Comparative et Performance")
        col_score_comp = 'Nombre des data set complétude >/=95%'
        col_score_prompt = nom_col_score_p
        df_synth = pd.merge(
            df_final_c[['Organisation unit', 'Complétude_Globale (%)', 'Reports_Actual', 'Reports_Attendu', col_score_comp]],
            df_final_p[['Organisation unit', 'Promptitude_Globale (%)', 'Reports_Actual_On_Time', col_score_prompt]],
            on='Organisation unit'
        )
        float_cols_synth = ['Complétude_Globale (%)', 'Promptitude_Globale (%)', 'Reports_Actual', 'Reports_Attendu', 'Reports_Actual_On_Time']
        for c in float_cols_synth:
            if c in df_synth.columns:
                df_synth[c] = pd.to_numeric(df_synth[c], errors='coerce').fillna(0).round(2)

        m1, m2, m3, m4, m5 = st.columns(5)
        m1.metric("Complétude moyenne", f"{df_synth['Complétude_Globale (%)'].mean():.2f}%")
        m2.metric("Promptitude moyenne", f"{df_synth['Promptitude_Globale (%)'].mean():.2f}%")
        m3.metric("Total Attendu", f"{df_synth['Reports_Attendu'].sum():.2f}")
        m4.metric("Complétude >= 95%", (df_final_c[col_rate_uniquement] >= 95).sum(axis=1).sum() if col_rate_uniquement else 0)
        m5.metric("Promptitude >= 95%", (df_final_p[col_rate_ot] >= 95).sum(axis=1).sum() if col_rate_ot else 0)

        st.subheader("📊 Tableau comparatif Complétude / Promptitude")
        df_comparatif = (
            df_synth[
                [
                    'Organisation unit',
                    'Complétude_Globale (%)',
                    col_score_comp,
                    'Promptitude_Globale (%)',
                    col_score_prompt
                ]
            ]
            .groupby('Organisation unit', as_index=False)
            .agg({
                'Complétude_Globale (%)': 'mean',
                col_score_comp: 'mean',
                'Promptitude_Globale (%)': 'mean',
                col_score_prompt: 'mean'
            })
            .round(2)
        )
        df_comparatif = df_comparatif.rename(columns={
            'Organisation unit': 'Aire de santé',
            col_score_comp: 'Nombre de dataset complétude >/= 95%',
            col_score_prompt: 'Nombre de dataset promptitude >/= 95%'
        })
        st.dataframe(
            df_comparatif.style.format({
                'Complétude_Globale (%)': '{:.2f}',
                'Promptitude_Globale (%)': '{:.2f}',
                'Nombre de dataset complétude >/= 95%': '{:.2f}',
                'Nombre de dataset promptitude >/= 95%': '{:.2f}'
            }),
            use_container_width=True,
            key="grid_comp_prompt"
        )

        # TABLEAU FUSIONNE (zone filtrée): reporting + actual + scores dataset
        if is_zone_focus:
            st.divider()
            st.subheader(f"📋 Tableau fusionné des indicateurs dataset : {selected_zone}")

            row_zone = target_df.iloc[0] if not target_df.empty else pd.Series(dtype=object)
            fusion_rows = []

            # Lignes reporting rate (complétude/promptitude en % + score >=95)
            for c_col in col_rate_uniquement:
                base_name = c_col.replace('Reporting rate', '').strip()
                p_col = next((p for p in col_rate_ot if base_name in p), None)
                compl_val = pd.to_numeric(row_zone.get(c_col, np.nan), errors='coerce')
                promp_val = pd.to_numeric(row_zone.get(p_col, np.nan), errors='coerce') if p_col else np.nan
                fusion_rows.append({
                    "Indicateurs (dataset)": f"{base_name} - Reporting",
                    "Complétude": round(float(compl_val), 2) if pd.notna(compl_val) else np.nan,
                    "Nombre de dataset complétude >/= 95%": int(compl_val >= 95) if pd.notna(compl_val) else np.nan,
                    "Promptitude": round(float(promp_val), 2) if pd.notna(promp_val) else np.nan,
                    "Nombre de dataset promptitude >/= 95%": int(promp_val >= 95) if pd.notna(promp_val) else np.nan
                })

            # Lignes actual reports (afficher tous les actual)
            for a_col in col_actual_strict:
                base_name = a_col.replace('Actual reports', '').strip()
                aot_col = next((a for a in actual_ot_cols if base_name in a), None)
                actual_val = pd.to_numeric(row_zone.get(a_col, np.nan), errors='coerce')
                actual_ot_val = pd.to_numeric(row_zone.get(aot_col, np.nan), errors='coerce') if aot_col else np.nan
                fusion_rows.append({
                    "Indicateurs (dataset)": f"{base_name} - Actual",
                    "Complétude": round(float(actual_val), 2) if pd.notna(actual_val) else np.nan,
                    "Nombre de dataset complétude >/= 95%": np.nan,
                    "Promptitude": round(float(actual_ot_val), 2) if pd.notna(actual_ot_val) else np.nan,
                    "Nombre de dataset promptitude >/= 95%": np.nan
                })

            df_tab4_fusion = pd.DataFrame(fusion_rows)
            st.dataframe(
                df_tab4_fusion.style.format({
                    "Complétude": "{:.2f}",
                    "Promptitude": "{:.2f}",
                    "Nombre de dataset complétude >/= 95%": "{:.0f}",
                    "Nombre de dataset promptitude >/= 95%": "{:.0f}"
                }),
                use_container_width=True,
                key="grid_fusion_zone"
            )

        st.subheader("Quadrant de Performance")
        fig_quad = px.scatter(df_synth, x='Complétude_Globale (%)', y='Promptitude_Globale (%)', text='Organisation unit', range_x=[0, 110], range_y=[0, 110])
        fig_quad.add_hline(y=80, line_dash="dot", line_color="red")
        fig_quad.add_vline(x=80, line_dash="dot", line_color="red")
        st.plotly_chart(fig_quad, use_container_width=True, key="chart_quad_tab4")

        col_l, col_r = st.columns(2)
        with col_l:
            st.success("🏆 **Top 5 Complétude**")
            df_top5_export = df_synth.nlargest(5, 'Complétude_Globale (%)')[['Organisation unit', 'Complétude_Globale (%)']].round(1)
            st.table(df_top5_export)
        with col_r:
            st.error("⚠️ **Flop 5 Promptitude**")
            df_flop5_export = df_synth.nsmallest(5, 'Promptitude_Globale (%)')[['Organisation unit', 'Promptitude_Globale (%)']].round(1)
            st.table(df_flop5_export)

    # --- ONGLET 5 : ÉLÉMENTS DE CATÉGORISATION (CORRECTIF APPLIQUÉ SUR TON SCRIPT) ---
    with tab5:
        st.header("🩺 Éléments de catégorisation interactifs")

        vr_groups = get_validation_groups(cache_user=current_user_for_cache)
        if vr_groups:
            group_mapping = {g['displayName']: g['id'] for g in vr_groups}
            options_regles = ["Toutes les règles (Global)"] + list(group_mapping.keys())

            selected_vr_name = st.selectbox("Sélectionner le Groupe de Règles de Validation :", options_regles)

            if st.button("Lancer l'analyse des violations"):
                can_run_validation = True
                if is_zone_focus and (not current_id_systeme or len(str(current_id_systeme)) <= 5):
                    can_run_validation = False

                if can_run_validation:
                    with st.spinner('Analyse des violations en cours...'):

                        target_group_id = group_mapping.get(selected_vr_name) if selected_vr_name != "Toutes les règles (Global)" else None

                        # 1) Construire la liste des zones selon le niveau de filtre
                        zone_id_mapping = {
                            str(k).strip(): str(v).strip()
                            for k, v in zone_to_id.items()
                            if str(k).strip() and str(v).strip()
                        }

                        if is_zone_focus:
                            zones_cibles = [selected_zone]
                            if selected_zone not in zone_id_mapping:
                                zid = mapping_ou_id.get(selected_zone)
                                if zid:
                                    zone_id_mapping[selected_zone] = str(zid)
                        else:
                            if filter_level == "Province" and selected_province != "Toutes les provinces":
                                zone_source_df = scope_df[
                                    (scope_df["Province"] == selected_province) &
                                    (scope_df["is_zone_candidate"])
                                ]
                            else:
                                zone_source_df = scope_df[scope_df["is_zone_candidate"]]

                            zones_cibles = sorted(zone_source_df["Zone"].dropna().unique().tolist())

                            for _, row_meta in zone_source_df.iterrows():
                                zone_name = str(row_meta.get("Zone", "")).strip()
                                zone_id = str(row_meta.get("Zone_ID", "")).strip()
                                if zone_name and zone_id:
                                    zone_id_mapping[zone_name] = zone_id

                            if not zones_cibles and current_id_systeme:
                                children_info = get_children_org_units_details(current_id_systeme, cache_user=current_user_for_cache)
                                zones_cibles = [c['displayName'] for c in children_info if c.get('displayName') and c.get('id')]
                                for child in children_info:
                                    zone_id_mapping[str(child['displayName']).strip()] = str(child['id']).strip()

                        zones_cibles = [z for z in zones_cibles if zone_id_mapping.get(z)]
                        zones_cibles = list(dict.fromkeys(zones_cibles))

                        # Mapping normalisé (fallback si les libellés diffèrent légèrement)
                        zone_id_mapping_norm = {
                            normalize_org_name(k): v
                            for k, v in zone_id_mapping.items()
                            if k and v
                        }

                        zones_resolues = []
                        for z in zones_cibles:
                            zid = zone_id_mapping.get(z) or zone_id_mapping_norm.get(normalize_org_name(z))
                            if zid:
                                zones_resolues.append((z, zid))
                        zones_cibles = [z for z, _ in zones_resolues]
                        zone_id_mapping = {z: zid for z, zid in zones_resolues}

                        if not zones_cibles:
                            st.warning("Aucune zone de santé valide trouvée.")
                            st.stop()

                        # 2) Recalcul des métriques au niveau "Zone de santé" (évite les zéros dus aux différences de niveaux)
                        actual_zone_by_id = {}
                        actual_zone_by_name = {}
                        comp_zone_by_id = {}
                        comp_zone_by_name = {}
                        prompt_zone_by_id = {}
                        prompt_zone_by_name = {}

                        scope_meta = scope_df[['Organisation unit', 'Organisation unit ID', 'Zone', 'Zone_ID']].copy()
                        scope_meta['Organisation unit'] = scope_meta['Organisation unit'].astype(str)
                        scope_meta['Organisation unit ID'] = scope_meta['Organisation unit ID'].astype(str)
                        scope_meta['Zone'] = scope_meta['Zone'].astype(str)
                        scope_meta['Zone_ID'] = scope_meta['Zone_ID'].astype(str)
                        scope_meta['_ou_key'] = scope_meta['Organisation unit'].map(normalize_org_name)
                        scope_meta = scope_meta.drop_duplicates(subset=['_ou_key', 'Organisation unit ID'])

                        def _attach_zone_meta(df_source):
                            if df_source is None or df_source.empty or 'Organisation unit' not in df_source.columns:
                                return pd.DataFrame()
                            tmp = df_source.copy()
                            tmp['Organisation unit'] = tmp['Organisation unit'].astype(str)
                            tmp['_ou_key'] = tmp['Organisation unit'].map(normalize_org_name)
                            tmp = tmp.merge(
                                scope_meta[['_ou_key', 'Organisation unit ID', 'Zone', 'Zone_ID']],
                                on='_ou_key',
                                how='left'
                            )
                            tmp['ou_id'] = tmp['Organisation unit'].map(mapping_ou_id)
                            tmp['ou_id'] = tmp['ou_id'].where(
                                tmp['ou_id'].notna() & (tmp['ou_id'].astype(str).str.strip() != ""),
                                tmp['Organisation unit ID']
                            )
                            tmp['Zone_ID'] = tmp['Zone_ID'].where(
                                tmp['Zone_ID'].notna() & (tmp['Zone_ID'].astype(str).str.strip() != ""),
                                tmp['ou_id']
                            )
                            tmp['Zone'] = tmp['Zone'].where(
                                tmp['Zone'].notna() & (tmp['Zone'].astype(str).str.strip() != ""),
                                tmp['Organisation unit']
                            )
                            return tmp

                        if {'Organisation unit', 'Reports_Actual', 'Reports_Attendu'}.issubset(df_synthese.columns):
                            df_comp_zone = _attach_zone_meta(
                                df_synthese[['Organisation unit', 'Reports_Actual', 'Reports_Attendu']]
                            )
                            if not df_comp_zone.empty:
                                df_comp_zone['Reports_Actual'] = pd.to_numeric(df_comp_zone['Reports_Actual'], errors='coerce').fillna(0)
                                df_comp_zone['Reports_Attendu'] = pd.to_numeric(df_comp_zone['Reports_Attendu'], errors='coerce').fillna(0)
                                grouped_comp = (
                                    df_comp_zone
                                    .groupby(['Zone_ID', 'Zone'], as_index=False)[['Reports_Actual', 'Reports_Attendu']]
                                    .sum()
                                )
                                grouped_comp['Complétude_Globale'] = (
                                    grouped_comp['Reports_Actual'] / grouped_comp['Reports_Attendu'].replace(0, np.nan) * 100
                                ).fillna(0).round(2)

                                actual_zone_by_id = grouped_comp.set_index('Zone_ID')['Reports_Actual'].to_dict()
                                comp_zone_by_id = grouped_comp.set_index('Zone_ID')['Complétude_Globale'].to_dict()
                                grouped_comp['_zone_key'] = grouped_comp['Zone'].map(normalize_org_name)
                                actual_zone_by_name = grouped_comp.set_index('_zone_key')['Reports_Actual'].to_dict()
                                comp_zone_by_name = grouped_comp.set_index('_zone_key')['Complétude_Globale'].to_dict()

                        if {'Organisation unit', 'Reports_Actual_On_Time', 'Reports_Attendu'}.issubset(df_final_p.columns):
                            df_prompt_zone = _attach_zone_meta(
                                df_final_p[['Organisation unit', 'Reports_Actual_On_Time', 'Reports_Attendu']]
                            )
                            if not df_prompt_zone.empty:
                                df_prompt_zone['Reports_Actual_On_Time'] = pd.to_numeric(df_prompt_zone['Reports_Actual_On_Time'], errors='coerce').fillna(0)
                                df_prompt_zone['Reports_Attendu'] = pd.to_numeric(df_prompt_zone['Reports_Attendu'], errors='coerce').fillna(0)
                                grouped_prompt = (
                                    df_prompt_zone
                                    .groupby(['Zone_ID', 'Zone'], as_index=False)[['Reports_Actual_On_Time', 'Reports_Attendu']]
                                    .sum()
                                )
                                grouped_prompt['Promptitude_Globale'] = (
                                    grouped_prompt['Reports_Actual_On_Time'] / grouped_prompt['Reports_Attendu'].replace(0, np.nan) * 100
                                ).fillna(0).round(2)

                                prompt_zone_by_id = grouped_prompt.set_index('Zone_ID')['Promptitude_Globale'].to_dict()
                                grouped_prompt['_zone_key'] = grouped_prompt['Zone'].map(normalize_org_name)
                                prompt_zone_by_name = grouped_prompt.set_index('_zone_key')['Promptitude_Globale'].to_dict()

                        # 3) Violations par zone + ratio /100 + corrigées T vs T-1
                        rows_cat = []
                        for zone_name in zones_cibles:
                            zone_id = zone_id_mapping.get(zone_name) or zone_id_mapping_norm.get(normalize_org_name(zone_name))
                            zone_key = normalize_org_name(zone_name)
                            comp_val = float(comp_zone_by_id.get(zone_id, comp_zone_by_name.get(zone_key, 0)))
                            prompt_val = float(prompt_zone_by_id.get(zone_id, prompt_zone_by_name.get(zone_key, 0)))
                            actual_val = float(actual_zone_by_id.get(zone_id, actual_zone_by_name.get(zone_key, 0)))
                            score_qualite = round((comp_val + prompt_val) / 2, 1)
                            if not zone_id or len(str(zone_id)) < 5:
                                rows_cat.append({
                                    'Zone de santé': zone_name,
                                    'Complétude_Globale': comp_val,
                                    'Promptitude_Globale': prompt_val,
                                    'Actual_Reports': actual_val,
                                    'Violations_Actuelles': 0,
                                    'Violations_Mois_Precedent': 0,
                                    'Violations_Corrigees': 0,
                                    'Score_Qualite': score_qualite
                                })
                                continue

                            df_val_zone = get_validation_results(zone_id, selection_mois, target_group_id, cache_user=current_user_for_cache)
                            df_val_prev_zone = get_validation_results(zone_id, selection_mois_prev, target_group_id, cache_user=current_user_for_cache)

                            # Comparaison des mêmes erreurs entre T-1 et T
                            sig_now = build_violation_signature_set(df_val_zone)
                            sig_prev = build_violation_signature_set(df_val_prev_zone)
                            violations_actuelles = len(sig_now)
                            violations_mois_precedent = len(sig_prev)
                            violations_corrigees = len(sig_prev - sig_now)

                            rows_cat.append({
                                'Zone de santé': zone_name,
                                'Complétude_Globale': comp_val,
                                'Promptitude_Globale': prompt_val,
                                'Actual_Reports': actual_val,
                                'Violations_Actuelles': int(violations_actuelles),
                                'Violations_Mois_Precedent': int(violations_mois_precedent),
                                'Violations_Corrigees': int(violations_corrigees),
                                'Score_Qualite': score_qualite
                            })

                        df_cat = pd.DataFrame(rows_cat)

                        # Formule demandée: nombre de règle violée / Actual * 100
                        df_cat['Ratio_Violations_sur_100'] = (
                            (df_cat['Violations_Actuelles'] / df_cat['Actual_Reports'].replace(0, np.nan)) * 100
                        ).fillna(0).round(2)

                        if 'Violations_Corrigees' not in df_cat.columns:
                            df_cat['Violations_Corrigees'] = 0
                        df_cat['Violations_Corrigees'] = pd.to_numeric(
                            df_cat['Violations_Corrigees'], errors='coerce'
                        ).fillna(0).astype(int)

                        st.subheader(f"Résultats : {selected_vr_name}")

                        df_display = df_cat[['Zone de santé', 'Complétude_Globale', 'Promptitude_Globale', 'Actual_Reports', 'Violations_Mois_Precedent', 'Violations_Corrigees', 'Violations_Actuelles', 'Ratio_Violations_sur_100', 'Score_Qualite']].rename(columns={
                            'Complétude_Globale': 'Complétude globale (%)',
                            'Promptitude_Globale': 'Promptitude globale (%)',
                            'Actual_Reports': 'Reports_Actual',
                            'Violations_Mois_Precedent': 'Règles violées (M-1)',
                            'Violations_Corrigees': 'Règles corrigées (M-1 -> M)',
                            'Violations_Actuelles': 'Règles violées (M)',
                            'Ratio_Violations_sur_100': 'Ratio / 100 rapports',
                            'Score_Qualite': 'Score de qualité'
                        })
                        df_tab5_export = df_display.copy()

                        st.dataframe(
                            df_display.style.format({
                                'Complétude globale (%)': '{:.2f}',
                                'Promptitude globale (%)': '{:.2f}',
                                'Reports_Actual': '{:.2f}',
                                'Ratio / 100 rapports': '{:.2f}',
                                'Score de qualité': '{:.2f}'
                            })
                            .map(style_taux, subset=['Complétude globale (%)', 'Promptitude globale (%)', 'Score de qualité'])
                            .highlight_max(subset=['Règles violées (M)'], color='#ffcccc')
                            .map(lambda x: 'color: red; font-weight: bold' if x > 10 else '', subset=['Ratio / 100 rapports']),
                            use_container_width=True
                        )

                        col_res1, col_res2 = st.columns(2)
                        col_res1.metric("Total violations détectées", int(df_cat['Violations_Actuelles'].sum()))
                        col_res2.metric("Total règles corrigées", int(df_cat['Violations_Corrigees'].sum()), delta=f"{int(df_cat['Violations_Corrigees'].sum())}")

                else:
                    st.error("L'ID de la zone de santé sélectionnée est manquant dans les données sources.")
        else:
            st.warning("⚠️ Impossible de charger les groupes de règles.")

    # --- MENU SIDEBAR : EXTRACTION COMMENTEE ---
    with st.sidebar:
        st.divider()
        st.subheader("📤 Extraction du rapport")
        report_type = st.selectbox(
            "Type de téléchargement :",
            ["Excel"],
            key="report_type_selector"
        )

        export_style_context = {
            "base": {
                "decimals": 2,
            },
            "actual_detail": {
                "decimals": 2,
            },
            "expected_detail": {
                "decimals": 2,
            },
            "completude": {
                "taux_cols": (col_rate_uniquement if 'col_rate_uniquement' in locals() else []) + ['Complétude_Globale (%)'],
                "score_cols": ['Nombre des data set complétude >/=95%'],
                "int_cols": ['Nombre des data set complétude >/=95%'],
                "decimals": 2,
            },
            "promptitude": {
                "taux_cols": (col_rate_ot if 'col_rate_ot' in locals() else []) + ['Promptitude_Globale (%)'],
                "score_cols": [nom_col_score_p] if 'nom_col_score_p' in locals() else [],
                "int_cols": [nom_col_score_p] if 'nom_col_score_p' in locals() else [],
                "decimals": 2,
            },
            "synthese": {
                "taux_cols": ['Complétude_Globale (%)', 'Promptitude_Globale (%)'],
                "decimals": 2,
            },
            "comparatif": {
                "taux_cols": ['Complétude_Globale (%)', 'Promptitude_Globale (%)'],
                "score_cols": ['Nombre de dataset complétude >/= 95%', 'Nombre de dataset promptitude >/= 95%'],
                "int_cols": ['Nombre de dataset complétude >/= 95%', 'Nombre de dataset promptitude >/= 95%'],
                "decimals": 2,
            },
            "fusion": {
                "taux_cols": ['Complétude', 'Promptitude'],
                "score_cols": ['Nombre de dataset complétude >/= 95%', 'Nombre de dataset promptitude >/= 95%'],
                "int_cols": ['Nombre de dataset complétude >/= 95%', 'Nombre de dataset promptitude >/= 95%'],
                "decimals": 2,
            },
            "tab5": {
                "taux_cols": ['Complétude globale (%)', 'Promptitude globale (%)', 'Score de qualité'],
                "highlight_max_cols": ['Règles violées (M)'],
                "ratio_alert_col": 'Ratio / 100 rapports',
                "ratio_alert_threshold": 10.0,
                "int_cols": ['Règles violées (M-1)', 'Règles corrigées (M-1 -> M)', 'Règles violées (M)'],
                "decimals": 2,
            },
            "top5": {
                "taux_cols": ['Complétude_Globale (%)'],
                "decimals": 2,
            },
            "flop5": {
                "taux_cols": ['Promptitude_Globale (%)'],
                "decimals": 2,
            },
        }

        comments_df = build_dashboard_comments_df(
            df_final_c=df_final_c,
            df_final_p=df_final_p,
            df_synth=df_synth,
            selected_zone=selected_scope_label,
            df_tab4_fusion=df_tab4_fusion,
            df_tab5=df_tab5_export,
            df_base=df_tab1_export,
            df_actual_detail=df_actual_detail_export,
            df_expected_detail=df_expected_detail_export,
            df_comparatif=df_comparatif,
            df_top5=df_top5_export,
            df_flop5=df_flop5_export
        )

        export_bytes = None
        export_mime = None
        export_name = None
        export_error = None

        export_bytes = to_excel_dashboard_report(
            df_base=df_tab1_export,
            df_actual_detail=df_actual_detail_export,
            df_expected_detail=df_expected_detail_export,
            df_final_c=df_final_c,
            df_final_p=df_final_p,
            df_synth=df_synth,
            comments_df=comments_df,
            df_tab4_fusion=df_tab4_fusion,
            df_comparatif=df_comparatif,
            df_tab5=df_tab5_export,
            df_top5=df_top5_export,
            df_flop5=df_flop5_export,
            style_context=export_style_context
        )
        export_mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        export_name = f"rapport_dashboard_commente_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx"

        if export_bytes is not None:
            st.download_button(
                f"Télécharger {report_type}",
                data=export_bytes,
                file_name=export_name,
                mime=export_mime,
                use_container_width=True,
                key=f"btn_download_report_{report_type}"
            )
        else:
            st.info(export_error if export_error else "Export indisponible.")

        st.divider()
        if "show_report_preview" not in st.session_state:
            st.session_state["show_report_preview"] = False
        col_prev1, col_prev2 = st.columns(2)
        if col_prev1.button("Visualiser le rapport", use_container_width=True):
            st.session_state["show_report_preview"] = True
        if col_prev2.button("Masquer la vue", use_container_width=True):
            st.session_state["show_report_preview"] = False

    if st.session_state.get("show_report_preview", False):
        st.divider()
        st.subheader("👁️ Visualisation consolidée du rapport")
        st.caption("Aperçu interactif des manipulations en cours (filtres, période, niveau).")

        p_tabs = st.tabs([
            "Base de données",
            "Données réelles",
            "Participants/Attendus",
            "Perf finale",
            "Promptitude",
            "Comparatif",
            "Top/Flop",
            "Règles",
            "Commentaires",
        ])

        with p_tabs[0]:
            st.dataframe(df_tab1_export, use_container_width=True)
        with p_tabs[1]:
            st.dataframe(df_actual_detail_export, use_container_width=True)
        with p_tabs[2]:
            st.dataframe(df_expected_detail_export, use_container_width=True)
        with p_tabs[3]:
            st.dataframe(df_final_c, use_container_width=True)
        with p_tabs[4]:
            st.dataframe(df_final_p, use_container_width=True)
        with p_tabs[5]:
            st.dataframe(df_comparatif, use_container_width=True)
            if not df_tab4_fusion.empty:
                st.write("Tableau fusionné zone filtrée")
                st.dataframe(df_tab4_fusion, use_container_width=True)
        with p_tabs[6]:
            cprev1, cprev2 = st.columns(2)
            with cprev1:
                st.write("Top 5 complétude")
                st.dataframe(df_top5_export, use_container_width=True)
            with cprev2:
                st.write("Flop 5 promptitude")
                st.dataframe(df_flop5_export, use_container_width=True)
        with p_tabs[7]:
            st.dataframe(df_tab5_export, use_container_width=True)
        with p_tabs[8]:
            st.dataframe(comments_df, use_container_width=True)
else:
    st.error("❌ Données indisponibles.")
